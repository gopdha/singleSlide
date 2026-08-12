"""
Slide Generation Agent — two genuinely different modes, not one component.

Mode 1 (run_slide_generation_discovery, once per project, at onboarding): a single
agentic call proposes 3 DESIGN PARAMETER sets (palette/fonts/flex_bounds) — never raw
pptx manipulation. Each set is rendered, deterministically, into one of 3 fixed layout
archetypes already implemented in this file. The PM picks one via ask_human (same
callable pattern as discovery_agent, same reasoning — see docs/DECISION_LOG.md). The
choice locks skills/<project_id>/slide-generation-agent/SKILL.md.

Mode 2 (render_report, every week after): fully deterministic, zero LLM calls, zero
MCP — a pure-ish function like core/rag_rollup.py. Loads the locked skill, renders the 3
report slides (Executive Summary, Feature Status, Initiative Status) into the locked
archetype, auto-fitting within the skill's flex_bounds (font size, then row height, then
truncation — Requirement 15's own lever order). If content doesn't fit even at the
tightest bounds, raises SlideFitError — Requirement 16's escalation, with no
review_gate/ built yet to hand it to; the caller inherits it, same as critique_agent's
risk-floor RuntimeError today.

WHY REQUIREMENT 17 ("Claude Agent SDK Skills: the built-in pptx skill...") IS SUPERSEDED:
Investigated the same way AskUserQuestion was investigated for discovery_agent (gotcha
#13) before writing any code here. Findings, in order of directness:
  1. claude_agent_sdk's own generated tool schema (sdk-tools.d.ts — auto-generated from
     the real CLI, "DO NOT MODIFY BY HAND") has no `Skill` tool definition at all, and
     zero occurrences of "pptx"/"docx"/"xlsx" anywhere in the file.
  2. ClaudeAgentOptions.skills (the real SDK mechanism for enabling skills) discovers
     skills from filesystem SKILL.md files (user/project settings + installed plugins) —
     not a magic built-in registry; its own docstring says so.
  3. A full search of every place this machine would discover a skill from (all 29
     installed plugin skills under ~/.claude/plugins/marketplaces/, plus this project's
     own nonexistent .claude/) turned up zero pptx/docx/xlsx skill anywhere.
No pptx Skill exists in this environment — not "reachable but unverified" like
AskUserQuestion was, genuinely absent, with no invocation path either. python-pptx
(already a dependency — mcp_servers/ppt_mcp uses it to PARSE .pptx; this module uses it
to WRITE one, same library, opposite direction) is used instead, matching this project's
deterministic-where-reproducibility-matters principle more literally than an LLM-driven
pptx Skill would have anyway. See docs/DECISION_LOG.md for the full writeup.

KNOWN LIMITATION, documented rather than hidden: python-pptx has no real text-layout
engine — there is no way to ask "how many lines will this text wrap to at this font size
in this box" and get PowerPoint's actual answer. The fit heuristic below (
_estimate_block_height_in) is a deterministic APPROXIMATION (a fixed average-character-
width formula), not exact rendering measurement. It is reproducible — the same content
always gets the same fit decision — but it is not pixel-exact. Good enough for this
project's purpose (a bounded, auditable auto-fit ladder), not claimed to be more than
that.
"""
import asyncio
import json as _json
import sys
from contextlib import aclosing
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable, Optional

from claude_agent_sdk import ClaudeAgentOptions, query
from claude_agent_sdk.types import ResultMessage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "common"))
from skill_loader import load_skill, write_skill  # noqa: E402

SKILL_NAME = "slide-generation-agent"

# =============================================================================
# Geometry — a standard 16:9 executive-deck canvas, shared by all 3 archetypes
# =============================================================================

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MARGIN_IN = 0.5
TITLE_BAR_IN = 1.1
SIDEBAR_WIDTH_IN = 3.2
GUTTER_IN = 0.3
TITLE_LINE_IN = 0.28  # vertical space a row's bold title line takes above its body text

DEFAULT_ROW_PADDING_IN = 0.12  # single default shared by every drawer (_draw_row_list,
                                # _draw_narrative_block, _draw_table_slide). Previously each
                                # drawer hardcoded its own slightly different default (0.14
                                # vs. 0.1) — that inconsistency is exactly what let a real
                                # invalid candidate through undetected: a candidate's
                                # row_height_in_min only had to beat ONE of several silently
                                # different defaults to look valid. One shared constant means
                                # one thing to validate against (see _validate_flex_bounds).

LAYOUT_ARCHETYPES = ["single_column_narrative", "two_column_metrics_sidebar", "banner_header_grid"]

STATUS_COLOR_KEYS = {
    "On Track": "green", "At Risk": "amber", "Blocked": "red", "Needs Human Review": "needs_review",
}
RAG_COLOR_KEYS = {"Green": "green", "Amber": "amber", "Red": "red", "Unknown": "needs_review"}

SUMMARY_CHARS_MULTIPLIER = 4  # the executive summary is a full paragraph, not a single
                               # row's display_text — reuses the same flex_bounds fields
                               # scaled up, rather than asking the skill for a second,
                               # narrower set of bounds nobody would meaningfully reason
                               # about differently at template-pick time


class SlideFitError(RuntimeError):
    """Raised when content doesn't fit even at the tightest flex bounds (Requirement 16:
    escalate to the PM, never silently degrade further). No review_gate/ exists yet to
    catch this — the caller (eventually core/orchestrator.py) inherits it, exactly the
    way critique_agent's risk-floor RuntimeError is inherited by run_pipeline today."""


# =============================================================================
# Fit heuristic — deterministic, approximate (see module docstring), the one
# genuinely non-agentic piece of Requirement 15's "auto-adjusts within flex bounds"
# =============================================================================

AVG_CHAR_WIDTH_FACTOR = 0.5  # heuristic: average glyph width ~= 0.5 * font size, a
                              # common approximation for proportional sans-serif fonts
LINE_SPACING_FACTOR = 1.2


def _estimate_block_height_in(text: str, font_size_pt: float, box_width_in: float) -> float:
    chars_per_line = max(1, int((box_width_in * 72) / (font_size_pt * AVG_CHAR_WIDTH_FACTOR)))
    lines = max(1, -(-len(text) // chars_per_line))  # ceil division, no text -> 1 line
    return lines * font_size_pt * LINE_SPACING_FACTOR / 72


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


@dataclass
class FitResult:
    font_size_pt: float
    row_padding_in: float
    texts: list[str]


def _fit_rows(
    texts: list[str], box_width_in: float, available_height_in: float,
    font_size_default: float, font_size_min: float,
    row_padding_default_in: float, row_padding_min_in: float,
    chars_default: int, chars_min: int,
    context: str, extra_per_row_in: float = 0.0,
) -> FitResult:
    """Deterministic bounded ladder: shrink font size, then row padding, then truncate
    text — in that exact order, matching Requirement 15's own listed lever order. Raises
    SlideFitError if nothing in the bounded search fits (Requirement 16).

    Guards its own flex_bounds inputs (default >= min for every lever) rather than
    looping zero times and raising a confusing "doesn't fit" for what's actually an
    invalid skill/candidate — a skill file is hand-editable text, so this is a real,
    not hypothetical, input to validate loudly."""
    if font_size_default < font_size_min:
        raise ValueError(f"{context}: font_size_default ({font_size_default}) is below font_size_pt_min ({font_size_min}) — invalid flex_bounds.")
    if row_padding_default_in < row_padding_min_in:
        raise ValueError(f"{context}: row_padding_default_in ({row_padding_default_in}) is below row_height_in_min ({row_padding_min_in}) — invalid flex_bounds.")
    if chars_default < chars_min:
        raise ValueError(f"{context}: display_text_max_chars_default ({chars_default}) is below display_text_max_chars_min ({chars_min}) — invalid flex_bounds.")

    font_size = font_size_default
    while font_size >= font_size_min:
        for row_padding in (row_padding_default_in, row_padding_min_in):
            for chars_cap in (chars_default, chars_min):
                truncated = [_truncate(t, chars_cap) for t in texts]
                total = sum(
                    _estimate_block_height_in(t, font_size, box_width_in) + row_padding + extra_per_row_in
                    for t in truncated
                )
                if total <= available_height_in:
                    return FitResult(font_size_pt=font_size, row_padding_in=row_padding, texts=truncated)
        font_size -= 1

    raise SlideFitError(
        f"{context}: content does not fit even at the tightest flex bounds "
        f"(font_size_pt_min={font_size_min}, row_height_in_min={row_padding_min_in}, "
        f"display_text_max_chars_min={chars_min}) for {len(texts)} row(s). "
        f"Escalate to the PM per CLAUDE.md Requirement 16 — no review_gate/ exists yet, "
        f"so the caller must handle this."
    )


# =============================================================================
# Drawing primitives — thin wrappers over python-pptx
# =============================================================================

def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _new_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _set_background(slide, hex_color: str) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _add_rect(slide, x: float, y: float, w: float, h: float, fill_hex: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _add_text(
    slide, x: float, y: float, w: float, h: float, text: str, font_name: str, size_pt: float,
    color_hex: str, bold: bool = False, italic: bool = False, align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex_to_rgb(color_hex)
    return box


# =============================================================================
# Design-derived helpers — turn a report + design dict into drawable content
# =============================================================================

def _status_hex(design: dict[str, Any], status_label: str) -> str:
    return design["palette"][STATUS_COLOR_KEYS.get(status_label, "needs_review")]


def _rag_hex(design: dict[str, Any], rag_status: str) -> str:
    return design["palette"][RAG_COLOR_KEYS.get(rag_status, "needs_review")]


def _status_counts(curated_features: list[dict[str, Any]]) -> dict[str, int]:
    counts = {"On Track": 0, "At Risk": 0, "Blocked": 0, "Needs Human Review": 0}
    for f in curated_features:
        counts[f["status_label"]] = counts.get(f["status_label"], 0) + 1
    return counts


def _feature_rows(curated_features: list[dict[str, Any]], design: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        {"title": f"{f['title']}  [{f['status_label']}]", "text": f["display_text"], "color_hex": _status_hex(design, f["status_label"])}
        for f in curated_features
    ]


def _initiative_rows(curated_initiatives: list[dict[str, Any]], design: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        {"title": i["title"], "text": i["display_text"], "color_hex": design["palette"]["brand_primary"]}
        for i in curated_initiatives
    ]


# =============================================================================
# Composite drawers — shared across archetypes
# =============================================================================

def _draw_row_list(slide, rows: list[dict[str, Any]], box_x: float, box_y: float, box_width: float,
                    box_height: float, design: dict[str, Any], context: str) -> Optional[FitResult]:
    """rows: [{"title": str, "text": str, "color_hex": Optional[str]}, ...]. Fits via
    _fit_rows, then draws a colored left bar + bold title + body text per row, stacked
    vertically. Empty rows list is a no-op (caller decides the empty-state message)."""
    if not rows:
        return None
    fonts, flex = design["fonts"], design["flex_bounds"]
    text_width = box_width - 0.35
    fit = _fit_rows(
        [r["text"] for r in rows], box_width_in=text_width, available_height_in=box_height,
        font_size_default=fonts["body_size_pt"], font_size_min=flex["font_size_pt_min"],
        row_padding_default_in=DEFAULT_ROW_PADDING_IN, row_padding_min_in=flex["row_height_in_min"],
        chars_default=flex["display_text_max_chars_default"], chars_min=flex["display_text_max_chars_min"],
        context=context, extra_per_row_in=TITLE_LINE_IN,
    )
    heading_size = min(fit.font_size_pt + 2, fonts["heading_size_pt"])
    y = box_y
    for row, text in zip(rows, fit.texts):
        block_h = _estimate_block_height_in(text, fit.font_size_pt, text_width) + fit.row_padding_in + TITLE_LINE_IN
        if row.get("color_hex"):
            _add_rect(slide, box_x, y, 0.09, block_h - fit.row_padding_in, row["color_hex"])
        text_x = box_x + 0.25
        _add_text(slide, text_x, y, box_width - 0.25, TITLE_LINE_IN, row["title"], fonts["heading_font"], heading_size, design["palette"]["text_primary"], bold=True)
        _add_text(slide, text_x, y + TITLE_LINE_IN, box_width - 0.25, block_h - TITLE_LINE_IN, text, fonts["body_font"], fit.font_size_pt, design["palette"]["text_secondary"])
        y += block_h
    return fit


def _draw_narrative_block(slide, text: str, box_x: float, box_y: float, box_width: float,
                           box_height: float, design: dict[str, Any], context: str) -> FitResult:
    fonts, flex = design["fonts"], design["flex_bounds"]
    fit = _fit_rows(
        [text], box_width_in=box_width, available_height_in=box_height,
        font_size_default=fonts["body_size_pt"], font_size_min=flex["font_size_pt_min"],
        row_padding_default_in=DEFAULT_ROW_PADDING_IN, row_padding_min_in=flex["row_height_in_min"],
        chars_default=flex["display_text_max_chars_default"] * SUMMARY_CHARS_MULTIPLIER,
        chars_min=flex["display_text_max_chars_min"] * SUMMARY_CHARS_MULTIPLIER,
        context=context,
    )
    _add_text(slide, box_x, box_y, box_width, box_height, fit.texts[0], fonts["body_font"], fit.font_size_pt, design["palette"]["text_primary"])
    return fit


def _draw_table_slide(
    slide, headers: list[str], rows: list[tuple], col_widths_in: list[float], text_col_index: int,
    box_x: float, box_y: float, box_width: float, box_height: float, design: dict[str, Any], context: str,
):
    """rows: tuples matching headers; the (possibly long) text at text_col_index drives
    the fit ladder, other columns are short fixed labels. Returns (fit, table) so callers
    can apply extra per-row styling (e.g. a status-colored first cell)."""
    if not rows:
        return None, None
    fonts, flex = design["fonts"], design["flex_bounds"]
    text_col_width = col_widths_in[text_col_index]
    fit = _fit_rows(
        [r[text_col_index] for r in rows], box_width_in=text_col_width - 0.2, available_height_in=box_height - 0.4,
        font_size_default=fonts["body_size_pt"], font_size_min=flex["font_size_pt_min"],
        row_padding_default_in=DEFAULT_ROW_PADDING_IN, row_padding_min_in=flex["row_height_in_min"],
        chars_default=flex["display_text_max_chars_default"], chars_min=flex["display_text_max_chars_min"],
        context=context, extra_per_row_in=0.1,
    )

    n_rows, n_cols = len(rows) + 1, len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, Inches(box_x), Inches(box_y), Inches(box_width), Inches(box_height)).table
    for j, w in enumerate(col_widths_in):
        table.columns[j].width = Inches(w)

    table.rows[0].height = Inches(0.4)
    palette = design["palette"]
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb(palette["brand_primary"])
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = _hex_to_rgb("#FFFFFF")

    for i, row in enumerate(rows):
        text = fit.texts[i]
        row_h = _estimate_block_height_in(text, fit.font_size_pt, text_col_width - 0.2) + fit.row_padding_in + 0.1
        table.rows[i + 1].height = Inches(row_h)
        for j, val in enumerate(row):
            cell_val = text if j == text_col_index else val
            cell = table.cell(i + 1, j)
            cell.text = str(cell_val)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(fit.font_size_pt)
            run.font.color.rgb = _hex_to_rgb(palette["text_primary"])

    return fit, table


def _draw_title_bar(slide, title: str, design: dict[str, Any], rag_status: Optional[str] = None) -> None:
    palette, fonts = design["palette"], design["fonts"]
    _add_rect(slide, 0, 0, SLIDE_WIDTH_IN, TITLE_BAR_IN, palette["brand_primary"])
    _add_text(slide, MARGIN_IN, 0.25, 8.5, 0.6, title, fonts["heading_font"], fonts["heading_size_pt"], "#FFFFFF", bold=True)
    if rag_status:
        badge_w = 1.6
        _add_rect(slide, SLIDE_WIDTH_IN - MARGIN_IN - badge_w, 0.3, badge_w, 0.5, _rag_hex(design, rag_status))
        _add_text(slide, SLIDE_WIDTH_IN - MARGIN_IN - badge_w, 0.3, badge_w, 0.5, rag_status, fonts["heading_font"], 16, "#FFFFFF", bold=True, align=PP_ALIGN.CENTER)


def _empty_state(slide, message: str, box_x: float, box_y: float, box_width: float, design: dict[str, Any]) -> None:
    fonts, palette = design["fonts"], design["palette"]
    _add_text(slide, box_x, box_y, box_width, 0.4, message, fonts["body_font"], fonts["body_size_pt"], palette["text_secondary"], italic=True)


# =============================================================================
# Archetype 1 — single_column_narrative: full-width, vertically stacked rows
# =============================================================================

def _draw_single_column_narrative(prs: Presentation, report: dict[str, Any], design: dict[str, Any], debug: bool = False) -> None:
    palette = design["palette"]
    content_x, content_width = MARGIN_IN, SLIDE_WIDTH_IN - 2 * MARGIN_IN

    slide = _new_slide(prs)
    _set_background(slide, palette["background"])
    _draw_title_bar(slide, "Executive Summary", design, rag_status=report["rag_status"])
    content_y = TITLE_BAR_IN + 0.3
    content_height = SLIDE_HEIGHT_IN - content_y - MARGIN_IN - (0.4 if report.get("trend_line") else 0)
    _draw_narrative_block(slide, report["executive_summary"], content_x, content_y, content_width, content_height, design, "Executive Summary body")
    if report.get("trend_line"):
        _add_text(slide, content_x, SLIDE_HEIGHT_IN - MARGIN_IN - 0.35, content_width, 0.35,
                   report["trend_line"], design["fonts"]["body_font"], design["fonts"]["body_size_pt"] - 1, palette["text_secondary"], italic=True)

    slide2 = _new_slide(prs)
    _set_background(slide2, palette["background"])
    _draw_title_bar(slide2, "Feature Status", design)
    content_y2 = TITLE_BAR_IN + 0.3
    content_h2 = SLIDE_HEIGHT_IN - content_y2 - MARGIN_IN
    if report["curated_features"]:
        _draw_row_list(slide2, _feature_rows(report["curated_features"], design), content_x, content_y2, content_width, content_h2, design, "Feature Status slide")
    else:
        _empty_state(slide2, "No committed Features reportable this week.", content_x, content_y2, content_width, design)

    slide3 = _new_slide(prs)
    _set_background(slide3, palette["background"])
    _draw_title_bar(slide3, "Initiative Status", design)
    content_y3 = TITLE_BAR_IN + 0.3
    content_h3 = SLIDE_HEIGHT_IN - content_y3 - MARGIN_IN
    if report["curated_initiatives"]:
        _draw_row_list(slide3, _initiative_rows(report["curated_initiatives"], design), content_x, content_y3, content_width, content_h3, design, "Initiative Status slide")
    else:
        _empty_state(slide3, "No other initiatives reported this week.", content_x, content_y3, content_width, design)

    if debug:
        print("[slide_generation_agent debug] drew single_column_narrative deck (3 slides)")


# =============================================================================
# Archetype 2 — two_column_metrics_sidebar: left status sidebar + right content
# =============================================================================

def _draw_sidebar(slide, report: dict[str, Any], design: dict[str, Any]) -> None:
    palette, fonts = design["palette"], design["fonts"]
    sidebar_y = TITLE_BAR_IN + 0.3

    _add_rect(slide, MARGIN_IN, sidebar_y, SIDEBAR_WIDTH_IN, 1.2, _rag_hex(design, report["rag_status"]))
    _add_text(slide, MARGIN_IN, sidebar_y + 0.15, SIDEBAR_WIDTH_IN, 0.5, report["rag_status"].upper(), fonts["heading_font"], 22, "#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN_IN, sidebar_y + 0.65, SIDEBAR_WIDTH_IN, 0.4, "Overall status", fonts["body_font"], 12, "#FFFFFF", align=PP_ALIGN.CENTER)

    counts = _status_counts(report["curated_features"])
    y = sidebar_y + 1.5
    for label in ("On Track", "At Risk", "Blocked", "Needs Human Review"):
        _add_rect(slide, MARGIN_IN, y, 0.25, 0.25, _status_hex(design, label))
        _add_text(slide, MARGIN_IN + 0.35, y - 0.03, SIDEBAR_WIDTH_IN - 0.35, 0.3, f"{label}: {counts.get(label, 0)}", fonts["body_font"], 13, palette["text_primary"])
        y += 0.4

    if report.get("trend_line"):
        _add_text(slide, MARGIN_IN, y + 0.15, SIDEBAR_WIDTH_IN, 0.8, report["trend_line"], fonts["body_font"], 11, palette["text_secondary"], italic=True)


def _draw_two_column_metrics_sidebar(prs: Presentation, report: dict[str, Any], design: dict[str, Any], debug: bool = False) -> None:
    palette = design["palette"]
    main_x = MARGIN_IN + SIDEBAR_WIDTH_IN + GUTTER_IN
    main_width = SLIDE_WIDTH_IN - main_x - MARGIN_IN
    main_y = TITLE_BAR_IN + 0.3
    main_height = SLIDE_HEIGHT_IN - main_y - MARGIN_IN

    slide = _new_slide(prs)
    _set_background(slide, palette["background"])
    _draw_title_bar(slide, "Executive Summary", design)
    _draw_sidebar(slide, report, design)
    _draw_narrative_block(slide, report["executive_summary"], main_x, main_y, main_width, main_height, design, "Executive Summary body")

    slide2 = _new_slide(prs)
    _set_background(slide2, palette["background"])
    _draw_title_bar(slide2, "Feature Status", design)
    _draw_sidebar(slide2, report, design)
    if report["curated_features"]:
        _draw_row_list(slide2, _feature_rows(report["curated_features"], design), main_x, main_y, main_width, main_height, design, "Feature Status slide")
    else:
        _empty_state(slide2, "No committed Features reportable this week.", main_x, main_y, main_width, design)

    slide3 = _new_slide(prs)
    _set_background(slide3, palette["background"])
    _draw_title_bar(slide3, "Initiative Status", design)
    _draw_sidebar(slide3, report, design)
    if report["curated_initiatives"]:
        _draw_row_list(slide3, _initiative_rows(report["curated_initiatives"], design), main_x, main_y, main_width, main_height, design, "Initiative Status slide")
    else:
        _empty_state(slide3, "No other initiatives reported this week.", main_x, main_y, main_width, design)

    if debug:
        print("[slide_generation_agent debug] drew two_column_metrics_sidebar deck (3 slides)")


# =============================================================================
# Archetype 3 — banner_header_grid: full-width banner + tabular grid content
# =============================================================================

def _draw_banner_header_grid(prs: Presentation, report: dict[str, Any], design: dict[str, Any], debug: bool = False) -> None:
    palette = design["palette"]
    content_x, content_width = MARGIN_IN, SLIDE_WIDTH_IN - 2 * MARGIN_IN
    content_y = TITLE_BAR_IN + 0.3
    content_height = SLIDE_HEIGHT_IN - content_y - MARGIN_IN

    slide = _new_slide(prs)
    _set_background(slide, palette["background"])
    _draw_title_bar(slide, "Executive Summary", design, rag_status=report["rag_status"])
    summary_height = content_height - (0.4 if report.get("trend_line") else 0)
    _draw_narrative_block(slide, report["executive_summary"], content_x, content_y, content_width, summary_height, design, "Executive Summary body")
    if report.get("trend_line"):
        _add_text(slide, content_x, SLIDE_HEIGHT_IN - MARGIN_IN - 0.35, content_width, 0.35,
                   report["trend_line"], design["fonts"]["body_font"], design["fonts"]["body_size_pt"] - 1, palette["text_secondary"], italic=True)

    slide2 = _new_slide(prs)
    _set_background(slide2, palette["background"])
    _draw_title_bar(slide2, "Feature Status", design)
    if report["curated_features"]:
        col_widths = [1.4, 3.4, content_width - 4.8]
        rows = [(f["status_label"], f["title"], f["display_text"]) for f in report["curated_features"]]
        fit, table = _draw_table_slide(slide2, ["Status", "Title", "Update"], rows, col_widths, 2, content_x, content_y, content_width, content_height, design, "Feature Status slide")
        for i, f in enumerate(report["curated_features"]):
            cell = table.cell(i + 1, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(_status_hex(design, f["status_label"]))
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = _hex_to_rgb("#FFFFFF")
    else:
        _empty_state(slide2, "No committed Features reportable this week.", content_x, content_y, content_width, design)

    slide3 = _new_slide(prs)
    _set_background(slide3, palette["background"])
    _draw_title_bar(slide3, "Initiative Status", design)
    if report["curated_initiatives"]:
        col_widths = [3.2, content_width - 3.2]
        rows = [(i["title"], i["display_text"]) for i in report["curated_initiatives"]]
        _draw_table_slide(slide3, ["Initiative", "Update"], rows, col_widths, 1, content_x, content_y, content_width, content_height, design, "Initiative Status slide")
    else:
        _empty_state(slide3, "No other initiatives reported this week.", content_x, content_y, content_width, design)

    if debug:
        print("[slide_generation_agent debug] drew banner_header_grid deck (3 slides)")


_ARCHETYPE_DRAW_FUNCTIONS: dict[str, Callable] = {
    "single_column_narrative": _draw_single_column_narrative,
    "two_column_metrics_sidebar": _draw_two_column_metrics_sidebar,
    "banner_header_grid": _draw_banner_header_grid,
}


# =============================================================================
# Mode 2 — render_report: fully deterministic, zero LLM, zero MCP
# =============================================================================

def _build_and_save(report: dict[str, Any], design: dict[str, Any], output_path: Path, debug: bool = False) -> Path:
    archetype = design["layout_archetype"]
    if archetype not in _ARCHETYPE_DRAW_FUNCTIONS:
        raise ValueError(f"Unknown layout_archetype '{archetype}' — must be one of {LAYOUT_ARCHETYPES}")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    _ARCHETYPE_DRAW_FUNCTIONS[archetype](prs, report, design, debug=debug)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    if debug:
        print(f"[slide_generation_agent debug] saved '{archetype}' deck -> {output_path}")
    return output_path


def render_report(
    report: dict[str, Any], project_id: str, skills_root: Optional[str] = None,
    output_path: Optional[Path] = None, debug: bool = False,
) -> Path:
    """Mode 2. Loads the locked skill (raises FileNotFoundError if the project hasn't
    run Mode 1 yet — same loud-failure convention as skill_loader.load_skill everywhere
    else), renders the 3 report slides into the locked archetype, auto-fitting within
    the skill's flex_bounds. No LLM call, no MCP client — a pure function of (report,
    design), same testability tier as core/rag_rollup.py.

    Raises SlideFitError if content doesn't fit even at the tightest bounds — see
    SlideFitError's docstring for what "escalate" concretely means today."""
    skill = load_skill(project_id, SKILL_NAME, skills_root)
    design = skill.frontmatter
    out = Path(output_path) if output_path else Path(__file__).parent / f"{project_id}_{report.get('week_of', 'report')}.pptx"
    return _build_and_save(report, design, out, debug=debug)


# =============================================================================
# ask_human — identical pattern and implementation to discovery_agent.py's, not
# cross-imported (established precedent — see docs/DECISION_LOG.md)
# =============================================================================

def default_ask_human(question: str, options: Optional[list[str]] = None) -> str:
    print(f"\n{question}")
    if options:
        for i, opt in enumerate(options, 1):
            print(f"  {i}. {opt}")
        while True:
            raw = input("> ").strip()
            if raw in options:
                return raw
            if raw.isdigit() and 1 <= int(raw) <= len(options):
                return options[int(raw) - 1]
            print(f"Please enter one of: {', '.join(options)} (or its number).")
    return input("> ").strip()


# =============================================================================
# Shared agentic-call plumbing (zero tools, structured output) — same pattern as
# synthesis_agent/critique_agent/discovery_agent, duplicated not cross-imported
# =============================================================================

async def _run_agentic_call(
    prompt: str, system_prompt: str, schema: dict[str, Any], model: str, max_turns: int, caller_name: str,
) -> dict[str, Any]:
    options = ClaudeAgentOptions(
        system_prompt=system_prompt, model=model, allowed_tools=[], permission_mode="dontAsk",
        max_turns=max_turns, output_format={"type": "json_schema", "schema": schema},
    )

    final_result: Optional[Any] = None
    try:
        async with aclosing(query(prompt=prompt, options=options)) as messages:
            async for message in messages:
                if isinstance(message, ResultMessage):
                    if message.is_error:
                        raise RuntimeError(message.result or f"{caller_name} run did not succeed: {message.subtype}")
                    final_result = message.structured_output if message.structured_output is not None else message.result
    except Exception as err:  # noqa: BLE001
        if "invalid api key" in str(err).lower() or "invalid x-api-key" in str(err).lower():
            raise RuntimeError(
                f"{caller_name} failed: invalid ANTHROPIC_API_KEY. "
                "Set a real key from https://console.anthropic.com/settings/keys and try again."
            ) from err
        raise RuntimeError(f"{caller_name} failed: {err}") from err

    if final_result is None:
        raise RuntimeError(f"{caller_name} produced no result.")

    return _json.loads(final_result) if isinstance(final_result, str) else final_result


# =============================================================================
# Mode 1 — run_slide_generation_discovery: agentic PARAMETERS, deterministic render
# =============================================================================

_PALETTE_KEYS = ["background", "text_primary", "text_secondary", "brand_primary", "green", "amber", "red", "needs_review"]

CANDIDATE_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "name": {"type": "string"},
        "rationale": {"type": "string"},
        "palette": {
            "type": "object",
            "properties": {k: {"type": "string", "description": "hex color, e.g. \"#0B3D91\""} for k in _PALETTE_KEYS},
            "required": _PALETTE_KEYS,
            "additionalProperties": False,
        },
        "fonts": {
            "type": "object",
            "properties": {
                "heading_font": {"type": "string"}, "body_font": {"type": "string"},
                "heading_size_pt": {"type": "integer"}, "body_size_pt": {"type": "integer"},
            },
            "required": ["heading_font", "body_font", "heading_size_pt", "body_size_pt"],
            "additionalProperties": False,
        },
        "flex_bounds": {
            "type": "object",
            "properties": {
                "font_size_pt_min": {"type": "integer", "description": "must be meaningfully below body_size_pt"},
                "row_height_in_min": {"type": "number"},
                "display_text_max_chars_default": {"type": "integer"},
                "display_text_max_chars_min": {"type": "integer", "description": "must be below display_text_max_chars_default"},
            },
            "required": ["font_size_pt_min", "row_height_in_min", "display_text_max_chars_default", "display_text_max_chars_min"],
            "additionalProperties": False,
        },
    },
    "required": ["name", "rationale", "palette", "fonts", "flex_bounds"],
    "additionalProperties": False,
}

CANDIDATE_SCHEMA = {
    "type": "object",
    "properties": {"candidates": {"type": "array", "minItems": 3, "maxItems": 3, "items": CANDIDATE_ITEM_SCHEMA}},
    "required": ["candidates"],
    "additionalProperties": False,
}

DESIGN_SYSTEM_PROMPT = f"""You are designing 3 slide-template candidates for a weekly executive
leadership report (Executive Summary, Feature Status, Initiative Status slides), for a PM to
choose one from at onboarding. You have no tools — produce your structured answer directly.

Each candidate you return will be rendered into a FIXED layout archetype, in this exact order —
candidate 1 always renders into archetype 1, and so on:
  1. single_column_narrative — a single full-width column, vertically stacked narrative rows
  2. two_column_metrics_sidebar — a left status/metrics sidebar + a right content column
  3. banner_header_grid — a full-width colored banner header + tabular grid content

Design each candidate's palette/fonts/flex_bounds to suit ITS archetype's character — a
banner/grid archetype often suits a more formal, high-contrast palette; a narrative archetype
often suits calmer, more editorial tones. Don't reuse near-identical palettes across all 3 — the
PM needs genuinely distinguishable options to form a real preference.

palette colors are hex strings (e.g. "#0B3D91"). green/amber/red/needs_review map directly to
the existing 4-level Feature status taxonomy (On Track/At Risk/Blocked/Needs Human Review) and
the RAG rollup (Green/Amber/Red/Unknown) — never invent a 5th status color or repurpose these
for anything else.

flex_bounds are the auto-fit floors used every week after this one-time choice is locked in:
- font_size_pt_min must be MEANINGFULLY below body_size_pt (leave real room to shrink — a gap
  of only 1-2pt defeats the point of having a floor at all).
- row_height_in_min is compared against a FIXED render-side constant, not a number you're free
  to invent: the renderer's own default row spacing is {DEFAULT_ROW_PADDING_IN} inches, and
  row_height_in_min MUST be strictly less than that, with real room to spare — not a value close
  to or above it. Target roughly half of {DEFAULT_ROW_PADDING_IN} inches or less (e.g. 0.04-0.08
  inches). A value at or above {DEFAULT_ROW_PADDING_IN} inches is invalid and will be rejected.
- display_text_max_chars_min must be meaningfully below display_text_max_chars_default.
These are validated in code after you respond — an invalid ordering (a floor at or above its
own default, or row_height_in_min at/above {DEFAULT_ROW_PADDING_IN} inches) is rejected, not
silently accepted.

Return your answer as JSON matching the provided schema exactly."""


def _validate_flex_bounds(candidate: dict[str, Any], context: str) -> None:
    """Validates a candidate's flex_bounds for internal self-consistency IMMEDIATELY after
    the agentic call returns, before any render is attempted — the same 'verify
    mechanically-checkable claims in code, never trust the agent's self-report' convention
    already named in docs/DECISION_LOG.md, just not applied here until a real live run
    caught the gap: a candidate's row_height_in_min (0.28) exceeded the code's own default
    row padding, and the only place that was ever checked was deep inside _fit_rows at
    render time — a ValueError mid-draw, not a clean rejection up front.

    Whole-batch failure is deliberate, not a candidate-by-candidate drop: unlike
    SlideFitError (a legitimate outcome even for a perfectly valid candidate — the model
    never sees the sample content it'll be rendered against, so it can't guarantee its
    bounds will fit unseen future content), an internally self-contradictory flex_bounds is
    content-independent and violates an explicit numeric constraint DESIGN_SYSTEM_PROMPT
    already states plainly. One candidate failing this check is evidence the model didn't
    reliably follow that constraint on this attempt at all — nothing guarantees the other
    two candidates are trustworthy either, they just didn't happen to trip the check.
    Silently dropping the bad one and presenting the "survivors" as if nothing was wrong
    would hide that from the PM. See docs/DECISION_LOG.md for the fuller writeup."""
    fonts, flex = candidate["fonts"], candidate["flex_bounds"]
    errors = []
    if flex["font_size_pt_min"] > fonts["body_size_pt"]:
        errors.append(f"font_size_pt_min ({flex['font_size_pt_min']}) exceeds body_size_pt ({fonts['body_size_pt']})")
    if flex["row_height_in_min"] > DEFAULT_ROW_PADDING_IN:
        errors.append(f"row_height_in_min ({flex['row_height_in_min']}) exceeds the render default row padding ({DEFAULT_ROW_PADDING_IN})")
    if flex["display_text_max_chars_min"] > flex["display_text_max_chars_default"]:
        errors.append(
            f"display_text_max_chars_min ({flex['display_text_max_chars_min']}) exceeds "
            f"display_text_max_chars_default ({flex['display_text_max_chars_default']})"
        )
    if errors:
        raise ValueError(f"{context}: invalid flex_bounds — " + "; ".join(errors))


async def _generate_candidates(project_id: str, model: str, debug: bool = False) -> list[dict[str, Any]]:
    result = await _run_agentic_call(
        prompt=f"Design 3 slide-template candidates for project '{project_id}'.",
        system_prompt=DESIGN_SYSTEM_PROMPT, schema=CANDIDATE_SCHEMA, model=model, max_turns=3,
        caller_name="_generate_candidates",
    )
    candidates = result["candidates"]
    for i, candidate in enumerate(candidates):
        candidate["layout_archetype"] = LAYOUT_ARCHETYPES[i]
        _validate_flex_bounds(candidate, context=f"Candidate {i + 1} ('{candidate['name']}')")
    if debug:
        print(f"[slide_generation_agent debug] generated {len(candidates)} candidate(s): {[c['name'] for c in candidates]}")
    return candidates


def _sample_report_fixture(project_id: str) -> dict[str, Any]:
    """A deliberately varied synthetic report — one of each status label, a long
    display_text near the truncation ceiling, a populated initiatives list — used to
    preview template candidates against real layout stress when Mode 1 isn't given a
    real report dict (e.g. no live run_pipeline() has happened for this project yet)."""
    return {
        "project_id": project_id, "week_of": "2026-08-16", "rag_status": "Amber",
        "executive_summary": (
            "Overall status is Amber this week. The dashboard aggregation feature remains On "
            "Track and is expected to ship on schedule. The ADO connector module has moved to At "
            "Risk due to a vendor sandbox access delay that is now blocking integration testing; "
            "escalation is in progress with the vendor's account team. The reporting pipeline "
            "revamp is Blocked pending a decision on story-point tracking conventions for a newly "
            "split team. One item needs human review before its status can be confidently reported."
        ),
        "trend_line": "Unchanged from last week's Amber rating; one new At Risk item.",
        "curated_features": [
            {"feature_id": 1, "title": "Build Agentic Dashboard", "status_label": "On Track",
             "display_text": "All planned Stories closed this sprint; on track for the committed ship date."},
            {"feature_id": 2, "title": "ADO Connector Module", "status_label": "At Risk",
             "display_text": (
                 "Vendor sandbox access has been delayed for the second consecutive week, now "
                 "blocking integration testing for the core connector path; escalation opened with "
                 "the vendor's account team and a response is expected by end of week, but no "
                 "confirmed resolution date exists yet."
             )},
            {"feature_id": 3, "title": "Reporting Pipeline Revamp", "status_label": "Blocked",
             "display_text": "Blocked on a decision about story-point tracking conventions for a newly split team."},
            {"feature_id": 4, "title": "Legacy Data Migration Cleanup", "status_label": "Needs Human Review",
             "display_text": "Evidence was ambiguous this week; investigation could not confidently assign a status."},
        ],
        "curated_initiatives": [
            {"title": "Internal Hackathon Prep", "display_text": "Logistics finalized; kickoff scheduled for next week."},
        ],
        "features": [], "initiatives": [], "prior_week": None,
    }


async def run_slide_generation_discovery(
    project_id: str,
    ask_human: Callable[..., str] = default_ask_human,
    sample_report: Optional[dict[str, Any]] = None,
    skills_root: Optional[str] = None,
    model: str = "claude-sonnet-5",
    output_dir: Optional[Path] = None,
    debug: bool = False,
) -> Path:
    """Mode 1. One agentic call per attempt proposes 3 design-parameter candidates
    (never raw pptx manipulation); each is rendered deterministically via the same
    _build_and_save Mode 2 uses, against a real or synthetic sample report; the PM picks
    one (or rejects all 3 for a fresh attempt) via ask_human. Persists the chosen
    candidate's parameters + the PM's stated reasoning as the locked skill."""
    sample = sample_report or _sample_report_fixture(project_id)
    out_dir = Path(output_dir) if output_dir else Path(__file__).parent

    while True:
        candidates = await _generate_candidates(project_id, model, debug=debug)
        rendered: list[tuple[dict[str, Any], str, Path]] = []
        for i, candidate in enumerate(candidates):
            path = out_dir / f"{project_id}_candidate_{i + 1}_{candidate['layout_archetype']}.pptx"
            try:
                _build_and_save(sample, candidate, path, debug=debug)
                rendered.append((candidate, candidate["layout_archetype"], path))
            except SlideFitError as err:
                print(f"\nCandidate {i + 1} ('{candidate['name']}') could not be rendered: {err}")

        if not rendered:
            retry = ask_human(
                "None of the 3 candidates could be rendered (flex bounds too tight for the sample "
                "content). Try again?", options=["Yes, regenerate", "No, cancel"],
            )
            if retry == "Yes, regenerate":
                continue
            raise RuntimeError("Slide-generation discovery cancelled: no candidate could be rendered.")

        listing = "\n".join(f"  {i + 1}. {c['name']} ({archetype}) -> {path}" for i, (c, archetype, path) in enumerate(rendered))
        option_labels = [str(i + 1) for i in range(len(rendered))] + ["Regenerate"]
        choice = ask_human(
            f"{len(rendered)} template candidate(s) rendered. Open each .pptx and compare:\n{listing}\n\n"
            "Which do you want to lock in?",
            options=option_labels,
        )
        if choice == "Regenerate":
            continue
        chosen_candidate, chosen_archetype, chosen_path = rendered[int(choice) - 1]
        break

    pm_notes = ask_human(f"Locking in '{chosen_candidate['name']}' ({chosen_archetype}). Any notes on why, for the record? (blank to skip)")

    frontmatter = {
        "project_id": project_id, "layout_archetype": chosen_archetype,
        "palette": chosen_candidate["palette"], "fonts": chosen_candidate["fonts"],
        "flex_bounds": chosen_candidate["flex_bounds"],
    }
    body = f"## Design rationale\n\n{chosen_candidate['rationale']}\n"
    if pm_notes.strip():
        body += f"\n### PM's stated reasoning at pick time\n{pm_notes.strip()}\n"

    final_confirm = ask_human("Persist this skill?", options=["Yes, persist", "No, cancel"])
    if final_confirm != "Yes, persist":
        raise RuntimeError("Slide-generation discovery cancelled by PM before persisting the slide-generation-agent skill.")

    return write_skill(project_id, SKILL_NAME, frontmatter, body, skills_root=skills_root)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Slide Generation Agent. Run this yourself; see README.md.")
    subparsers = parser.add_subparsers(dest="command", required=True)

    discover = subparsers.add_parser("discover", help="Mode 1: interactive template selection (run once per project).")
    discover.add_argument("project_id")
    discover.add_argument("--sample-report-json", default=None, help="Path to a real report JSON to preview candidates against (defaults to a built-in synthetic fixture).")
    discover.add_argument("--output-dir", default=None)
    discover.add_argument("--skills-root", default=None)
    discover.add_argument("--debug", action="store_true")

    render = subparsers.add_parser("render", help="Mode 2: render a finalized report into the locked template.")
    render.add_argument("project_id")
    render.add_argument("report_json", help="Path to a finalized report JSON (synthesize_report()'s output shape).")
    render.add_argument("--output", default=None)
    render.add_argument("--skills-root", default=None)
    render.add_argument("--debug", action="store_true")

    args = parser.parse_args()

    async def _main():
        if args.command == "discover":
            sample = None
            if args.sample_report_json:
                sample = _json.loads(Path(args.sample_report_json).read_text(encoding="utf-8"))
            path = await run_slide_generation_discovery(
                args.project_id, sample_report=sample, skills_root=args.skills_root,
                output_dir=Path(args.output_dir) if args.output_dir else None, debug=args.debug,
            )
            print(f"\nDone. slide-generation-agent skill written to: {path}")
        else:
            report = _json.loads(Path(args.report_json).read_text(encoding="utf-8"))
            path = render_report(
                report, args.project_id, skills_root=args.skills_root,
                output_path=Path(args.output) if args.output else None, debug=args.debug,
            )
            print(f"\nDone. Rendered deck: {path}")

    asyncio.run(_main())
