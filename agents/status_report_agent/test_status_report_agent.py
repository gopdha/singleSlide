#!/usr/bin/env python3
"""
Structural + skill-loading checks always run. Request-construction validated with a
placeholder ANTHROPIC_API_KEY if set (using a real, minimal .pptx fixture so the MCP call
itself succeeds and the only expected failure is Anthropic auth). A full live run needs a
real API key + a directory of real submitted reports — see README.md.

Run with:
    python3 test_status_report_agent.py
    ANTHROPIC_API_KEY=sk-ant-fake python3 test_status_report_agent.py
    ANTHROPIC_API_KEY=... STATUS_REPORTS_DIR=... [ADO_ORG=... ADO_PROJECT=... ADO_PAT_BASE64=...] python3 test_status_report_agent.py
"""
import asyncio
import json
import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from status_report_agent import (  # noqa: E402
    STATUS_REPORT_SCHEMA,
    investigate_status_report,
    list_submitted_reports,
)

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "common"))
from skill_loader import load_skill  # noqa: E402

SKILLS_ROOT = str(Path(__file__).parent.parent.parent / "skills")
TEST_PROJECT_ID = "ai-reports-demo"
FIXTURE_PATH = Path(__file__).parent / "test-fixture-report.pptx"


def _build_fixture():
    """A minimal real .pptx so parse_slide succeeds — proves the MCP wiring and request
    construction are correct, isolating any failure to Anthropic auth (tier 2) or letting
    a real investigation run (tier 3, via a caller-provided real reports directory)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def add_text(text, top, size=14, bold=False):
        box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.4))
        tf = box.text_frame
        tf.text = text
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold

    add_text("Prepared by: Test Lead", Inches(0.2), size=11)
    add_text("Accomplishments", Inches(0.7), size=16, bold=True)
    add_text("Finished the structural test fixture.", Inches(1.1))
    add_text("Risks", Inches(1.6), size=16, bold=True)
    add_text("None currently.", Inches(2.0))

    prs.save(str(FIXTURE_PATH))


async def main():
    print("Status Report Agent (skill-driven) — test suite\n")

    print("Structural checks (no API/network call):")
    assert callable(investigate_status_report)
    assert callable(list_submitted_reports)
    print("  PASS  investigate_status_report and list_submitted_reports are importable and callable")

    other_initiative_props = set(STATUS_REPORT_SCHEMA["properties"]["other_initiatives"]["items"]["properties"].keys())
    assert "status_label" not in other_initiative_props, "Other Initiatives must stay narrative-only — no status_label, by design (see DECISION_LOG.md)"
    print("  PASS  Other Initiatives schema has no status_label field (narrative-only, as designed)")

    enrichment_props = set(STATUS_REPORT_SCHEMA["properties"]["potential_enrichments"]["items"]["properties"].keys())
    assert enrichment_props.issuperset({"related_feature_id", "related_feature_title_guess", "match_confidence", "excerpt", "note"})
    assert not enrichment_props & {"already_tracked", "duplicate_of_ado", "merge_decision", "already_in_evidence"}, \
        "potential_enrichments must only flag candidates — the merge decision belongs to Synthesis Agent, which has full cross-source visibility"
    assert STATUS_REPORT_SCHEMA["properties"]["potential_enrichments"]["items"]["properties"]["match_confidence"]["enum"] == ["high", "low"]
    print("  PASS  potential_enrichments schema only flags candidates (no merge-decision field) and carries match_confidence")

    print("\nSkill loading checks:")
    skill = load_skill(TEST_PROJECT_ID, "status-report-agent", SKILLS_ROOT)
    assert skill.frontmatter["report_file_glob"] == "*.pptx"
    assert "match_confidence" in skill.body or "confidence" in skill.body.lower()
    print("  PASS  example skill loads correctly (frontmatter + body)")

    try:
        load_skill("nonexistent-project", "status-report-agent", SKILLS_ROOT)
        print("  FAIL  expected FileNotFoundError for a project with no skill, got none")
        sys.exit(1)
    except FileNotFoundError:
        print("  PASS  missing skill fails loudly (no silent fallback to a default)")

    print("\nDeterministic listing check:")
    reports = list_submitted_reports(str(Path(__file__).parent), TEST_PROJECT_ID, SKILLS_ROOT)
    assert isinstance(reports, list)
    print(f"  PASS  list_submitted_reports() runs against a real directory (found {len(reports)} existing .pptx here)")

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("\nANTHROPIC_API_KEY not set — skipping request validation and live run.")
        print("Structural checks passed.")
        return

    _build_fixture()
    try:
        reports_dir = os.environ.get("STATUS_REPORTS_DIR")
        if not reports_dir:
            print("\nANTHROPIC_API_KEY is set but STATUS_REPORTS_DIR is not.")
            print("Running a request-construction check only, against a real fixture .pptx")
            print("(expects a clean failure at Anthropic auth, since parse_slide itself should succeed):\n")
            try:
                await investigate_status_report(
                    file_path=str(FIXTURE_PATH), existing_features=[{"id": 1, "title": "Placeholder Feature"}],
                    project_id=TEST_PROJECT_ID, skills_root=SKILLS_ROOT, model="claude-sonnet-5",
                )
                print("Unexpected: call succeeded with a placeholder key — please double-check this result.")
            except RuntimeError as err:
                print(f"Got expected failure (this is success for a structural check): {err}")
                print("\nRequest construction validated. Structural checks passed.")
            return

        existing_features = []
        existing_features_json = os.environ.get("EXISTING_FEATURES_JSON")
        org, project, pat_base64 = os.environ.get("ADO_ORG"), os.environ.get("ADO_PROJECT"), os.environ.get("ADO_PAT_BASE64")
        if existing_features_json:
            # A synthetic override for exercising enrichment-matching without needing real ADO
            # credentials — e.g. '[{"id": 101, "title": "Build Agentic Dashboard"}]'. Real runs
            # should prefer ADO_ORG/ADO_PROJECT/ADO_PAT_BASE64 below.
            existing_features = json.loads(existing_features_json)
            print(f"\nUsing {len(existing_features)} synthetic existing Feature(s) from EXISTING_FEATURES_JSON for enrichment matching.")
        elif org and project and pat_base64:
            sys.path.insert(0, str(Path(__file__).parent.parent / "feature_agent"))
            from feature_agent import list_committed_features  # noqa: E402
            existing_features = await list_committed_features(org, project, pat_base64, TEST_PROJECT_ID, SKILLS_ROOT)
            print(f"\nLoaded {len(existing_features)} existing Feature(s) from ADO for enrichment matching.")
        else:
            print("\nADO_ORG/ADO_PROJECT/ADO_PAT_BASE64 (and EXISTING_FEATURES_JSON) not set — running without "
                  "existing-Feature context (everything will land as an Other Initiative or a low-confidence guess).")

        reports = list_submitted_reports(reports_dir, TEST_PROJECT_ID, SKILLS_ROOT, debug=True)
        print(f"Found {len(reports)} submitted report(s): {[r['file_name'] for r in reports]}")
        if not reports:
            print("No reports found — drop at least one .pptx into STATUS_REPORTS_DIR and retry.")
            return

        first = reports[0]
        print(f"\nInvestigating {first['file_name']}...")
        result = await investigate_status_report(first["file_path"], existing_features, TEST_PROJECT_ID, SKILLS_ROOT)
        print("\nResult:")
        print(json.dumps(result, indent=2))

        assert "other_initiatives" in result and "potential_enrichments" in result
        for oi in result["other_initiatives"]:
            assert "status_label" not in oi
        print("\nPASS  live Status Report Agent investigation completed; Other Initiatives stayed label-free")
    finally:
        FIXTURE_PATH.unlink(missing_ok=True)


if __name__ == "__main__":
    asyncio.run(main())
