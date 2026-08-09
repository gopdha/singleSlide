#!/usr/bin/env python3
"""
Local test harness for the PPT MCP server.
Builds a small real .pptx test fixture (a plausible team-lead status slide with
text sections, a table, and color-coded shapes) purely to prove the MCP server
and parser work end-to-end over the real MCP protocol — not part of the
pipeline's actual weekly data.
"""
import asyncio
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from mcp import ClientSession
from mcp.client.stdio import StdioServerParameters, stdio_client

FIXTURE_PATH = Path(__file__).parent / "test-fixture-status.pptx"


def build_fixture():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    def add_text(text, top, left=Inches(0.5), width=Inches(9), size=14, bold=False):
        box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = box.text_frame
        tf.text = text
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        return box

    add_text("Prepared by: Priya Nair", Inches(0.2), size=11)

    add_text("Accomplishments", Inches(0.7), size=16, bold=True)
    add_text("Completed sprint 14 with all planned stories closed.", Inches(1.1))
    add_text("Shipped the new onboarding flow to staging.", Inches(1.45))

    add_text("Planned Next Steps", Inches(1.9), size=16, bold=True)
    add_text("Begin integration testing for the payments module.", Inches(2.3))

    add_text("Risks", Inches(2.75), size=16, bold=True)
    add_text("Vendor API sandbox has been unstable this week.", Inches(3.15))

    add_text("Help Needed", Inches(3.6), size=16, bold=True)
    add_text("Need a decision on the fallback vendor by Friday.", Inches(4.0))

    # A color-coded status shape (Amber) — simulates a RAG indicator common in these decks.
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(0.2), Inches(0.6), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)  # Amber
    shape.text_frame.text = ""

    # A small table (common for a mini task list).
    rows, cols = 3, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(4.5), Inches(6), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Task"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "Payments integration"
    table.cell(1, 1).text = "In progress"
    table.cell(2, 0).text = "Vendor eval"
    table.cell(2, 1).text = "Blocked"

    prs.save(str(FIXTURE_PATH))
    print(f"Built test fixture at {FIXTURE_PATH}")


async def main():
    build_fixture()

    server_params = StdioServerParameters(command=sys.executable, args=["server.py"], cwd=str(Path(__file__).parent))

    async with stdio_client(server_params) as (read, write):
        async with ClientSession(read, write) as session:
            await session.initialize()

            tools = await session.list_tools()
            print("\nTools exposed by server:", [t.name for t in tools.tools])

            result = await session.call_tool("parse_slide", {"file_path": str(FIXTURE_PATH)})
            print("\nparse_slide() result:")
            for block in result.content:
                if hasattr(block, "text"):
                    parsed = json.loads(block.text)
                    print(json.dumps(parsed, indent=2))

    FIXTURE_PATH.unlink(missing_ok=True)
    print("\nTest fixture cleaned up. Done.")


if __name__ == "__main__":
    asyncio.run(main())
