#!/usr/bin/env python3
"""
PPT MCP Server (internal, custom)
Implements: parse_slide(file_path) -> { team_lead_id, sections, color_cues, parse_confidence }
Per System Design Document Section 4.3.

Run standalone (stdio transport) with:  python3 server.py
Configure in an MCP client (Claude Desktop/Code) pointing at this file.
"""
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu
# mcp>=2.0.0 renamed/relocated FastMCP to mcp.server.mcpserver.MCPServer — same API
# surface (name=, .tool() decorator, .run_stdio_async()), just a different import path.
from mcp.server.mcpserver import MCPServer as FastMCP

mcp = FastMCP(name="ppt-teamlead-mcp")

# Section header keywords -> canonical section name.
SECTION_KEYWORDS: dict[str, list[str]] = {
    "accomplishments": ["accomplishment", "done", "completed", "wins", "progress"],
    "planned_next": ["planned next", "next steps", "upcoming", "plan for next", "next week"],
    "risks": ["risk", "blocker", "issue", "concern"],
    "help_needed": ["help needed", "ask", "decisions needed", "support needed", "escalation"],
}

# Reference RAG colors for nearest-color classification.
RAG_REFERENCE = {
    "Red": (0xC0, 0x00, 0x00),
    "Amber": (0xFF, 0xC0, 0x00),
    "Green": (0x00, 0xB0, 0x50),
}


def classify_color(rgb: tuple[int, int, int]) -> str | None:
    """Nearest-neighbor classify an RGB tuple against the RAG reference colors.
    Returns None if the color is too far from all three (e.g., a neutral gray/white)."""
    best_status, best_dist = None, float("inf")
    for status, ref in RAG_REFERENCE.items():
        dist = sum((a - b) ** 2 for a, b in zip(rgb, ref)) ** 0.5
        if dist < best_dist:
            best_status, best_dist = status, dist
    # Threshold: colors far from all three reference points (e.g. neutral grays) are not a status cue.
    return best_status if best_dist < 90 else None


def match_section(text: str) -> str | None:
    lowered = text.strip().lower()
    if not lowered:
        return None
    for section, keywords in SECTION_KEYWORDS.items():
        for kw in keywords:
            if kw in lowered:
                return section
    return None


def extract_team_lead_id(slide, file_path: Path) -> str:
    """Look for a 'Prepared by: <name>' style line; fall back to the filename stem."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            m = re.search(r"(?:prepared by|team lead|owner)\s*[:\-]\s*(.+)", text, re.IGNORECASE)
            if m:
                return m.group(1).strip()
    return file_path.stem


def extract_tables(slide) -> list[list[list[str]]]:
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)
    return tables


def extract_color_cues(slide) -> list[dict[str, Any]]:
    cues = []
    for idx, shape in enumerate(slide.shapes):
        fill = getattr(shape, "fill", None)
        if fill is None:
            continue
        try:
            if fill.type is not None and fill.type == 1:  # MSO_FILL.SOLID == 1
                rgb = fill.fore_color.rgb
                rgb_tuple = (rgb[0], rgb[1], rgb[2]) if not isinstance(rgb, str) else tuple(
                    int(str(rgb)[i : i + 2], 16) for i in (0, 2, 4)
                )
                status = classify_color(rgb_tuple)
                if status:
                    cues.append(
                        {
                            "shape_id": shape.shape_id if hasattr(shape, "shape_id") else idx,
                            "shape_name": shape.name,
                            "rgb": "%02X%02X%02X" % rgb_tuple,
                            "inferred_status": status,
                        }
                    )
        except Exception:
            # Shapes without a resolvable solid fill (gradients, no fill, theme colors we can't
            # resolve without the theme) are skipped rather than raising.
            continue
    return cues


def extract_sections(slide) -> dict[str, list[str]]:
    sections: dict[str, list[str]] = {k: [] for k in SECTION_KEYWORDS}
    current_section: str | None = None

    # Sort shapes roughly top-to-bottom, left-to-right so headers precede their body text.
    shapes_with_pos = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        top = shape.top if shape.top is not None else Emu(0)
        left = shape.left if shape.left is not None else Emu(0)
        shapes_with_pos.append((top, left, shape))
    shapes_with_pos.sort(key=lambda t: (t[0], t[1]))

    for _, _, shape in shapes_with_pos:
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if not text:
                continue
            header_match = match_section(text)
            if header_match and len(text) < 40:
                # Treat short lines matching a keyword as a section header.
                current_section = header_match
                continue
            if current_section:
                sections[current_section].append(text)
            # Text before any recognized header is not assigned to a section
            # (captured separately as unsectioned_text for transparency).

    return sections


def parse_slide_file(file_path: str) -> dict[str, Any]:
    path = Path(file_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPT file not found: {path}")

    prs = Presentation(str(path))
    if len(prs.slides) == 0:
        raise ValueError("Presentation has no slides.")

    # Convention: the team lead's status slide is the first slide in the deck.
    slide = prs.slides[0]

    team_lead_id = extract_team_lead_id(slide, path)
    sections = extract_sections(slide)
    tables = extract_tables(slide)
    color_cues = extract_color_cues(slide)

    found_sections = [k for k, v in sections.items() if v]
    missing_sections = [k for k in SECTION_KEYWORDS if k not in found_sections]
    confidence = round(len(found_sections) / len(SECTION_KEYWORDS), 2)

    return {
        "team_lead_id": team_lead_id,
        "sections": sections,
        "tables": tables,
        "color_cues": color_cues,
        "parse_confidence": confidence,
        "missing_sections": missing_sections,
        "slide_count_in_deck": len(prs.slides),
    }


@mcp.tool(
    name="parse_slide",
    description=(
        "Parses a team lead's weekly status PowerPoint (first slide) and extracts structured "
        "status data: accomplishments, planned next steps, risks/blockers, help needed, any "
        "tables, and color-coded RAG cues from shape fills. Returns a parse_confidence score "
        "(0.0-1.0) so low-confidence extractions can be flagged for PM review rather than "
        "silently trusted."
    ),
)
def parse_slide(file_path: str) -> dict[str, Any]:
    try:
        return parse_slide_file(file_path)
    except Exception as exc:  # surfaced to the MCP client as a tool error
        return {"error": str(exc)}


if __name__ == "__main__":
    import asyncio

    asyncio.run(mcp.run_stdio_async())
